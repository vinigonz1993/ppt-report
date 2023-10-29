from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import (
    XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
)
from pptx.util import Inches
from sample.samples import sample


class PPTReport:
    '''
        Class to manage the creation of a PPT presentation.
        It is possible to easily create a presentation using
        a dictionary of slides and slide properties.
    '''

    def __init__(self, name, properties):
        self.presentation = Presentation()
        self.layout = self.presentation.slide_layouts[0]
        self.name = name
        self.properties = properties

    def add_slide(self, properties):
        '''
            Adds a slide
        '''
        slide = self.presentation.slides.add_slide(self.layout)

        if "title" in properties:
            title = slide.shapes.title

            for attr in properties["title"]:
                setattr(title, attr, properties["title"][attr])

        if "text" in properties:
            text = slide.placeholders[1]

            for attr in properties["text"]:
                setattr(text, attr, properties["text"][attr])

        if "bar_chart" in properties:
            bar_chart = properties["bar_chart"]
            self.generate_chart(
                slide,
                bar_chart.get("title"),
                bar_chart.get("data", []),
                bar_chart.get("categories", []),
                bar_chart.get("labels", False),
                bar_chart.get("legend", False),
                bar_chart.get("chart_type", "BAR")
            )

        return slide

    def validate_chart_properties(self, props):
        '''Checks for invalid properties'''
        invalid = False
        if not props["data"]:
            print("Categories are missing")
            invalid = True

        if not props["categories"]:
            print("Categories are missing")
            invalid = True

        if invalid:
            quit()

    def chart_type(self, chart_type):
        types = {
            "BAR": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "LINE": XL_CHART_TYPE.LINE,
            "PIE": XL_CHART_TYPE.PIE
        }
        return types[chart_type]

    def generate_chart(
        self, slide, chart_title="", data=[],
        categories=[], labels=False, legend=False,
        chart_type="BAR"
    ):
        '''
            Adds a chart to the specific slide
        '''

        self.validate_chart_properties({
            "data": data,
            "categories": categories
        })

        chart_data = ChartData()
        chart_data.title = chart_title
        chart_data.categories = categories

        for serie in data:
            chart_data.add_series(
                serie[0],
                serie[1]
            )

        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
        chart = slide.shapes.add_chart(
            self.chart_type(chart_type), x, y, cx, cy, chart_data
        ).chart

        plot = chart.plots[0]
        plot.has_data_labels = labels

        if legend:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.series[0].smooth = True

        if labels and chart_type != "LINE":
            data_labels = plot.data_labels
            data_labels.position = XL_LABEL_POSITION.INSIDE_END

    def mount(self):
        '''Mounts the presentation'''
        for i in self.properties.keys():
            report.add_slide(properties=self.properties[i])

        self.presentation.save(f"{self.name}.pptx")


report = PPTReport("report", sample())

report.mount()
